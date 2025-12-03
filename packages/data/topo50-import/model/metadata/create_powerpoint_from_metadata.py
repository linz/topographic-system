import os
from pptx import Presentation
from pptx.util import Inches, Pt
import json

def add_image(slide, image_path, left=Inches(5), top=Inches(1), width=Inches(3), height=Inches(7.5)):
    if image_path and os.path.exists(image_path):
        if left == -1:
            slide.shapes.add_picture(image_path, 1,3)
        else:
            slide.shapes.add_picture(image_path, left, top, width=width, height=height)

def add_summary_slide(prs, summary_info):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    category = summary_info[0].title()
    if title_shape:
        title_shape.text = f"{category} Layers and Themes"
    else:
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = f"{category} Layers and Themes"

    image_path = summary_info[1].replace('.svg','.png')
    add_image(slide, image_path, left=-1)

def add_layer_slide(prs, layer_name, feature_names, layer_description, model_picture_path, layer_category=''):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    layer_name = layer_name.replace('_', ' ').title()
    if title_shape:
        title_shape.text = f"Layer: {layer_name}"
    else:
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = f"Layer: {layer_name}"
    
    # Add feature names list
    feature_list = "\n".join([f"• {feature}" for feature in feature_names])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(3))
    tf = txBox.text_frame
    tf.text = f"Layer Category: {layer_category}\nFeatures:\n{feature_list}\n\nLayer Description:\n{layer_description}"
    
    add_image(slide, model_picture_path, left=Inches(8.6), top=Inches(0), width=Inches(3.5), height=Inches(7.5))

def add_feature_slide(prs, feature_type, theme, table_lds, table_type, description, feature_picture_path, aerial_picture_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    feature_type = feature_type.replace('_', ' ').title()
    if title_shape:
        title_shape.text = f"Feature Type: {feature_type}"
    else:
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = f"Feature Type: {feature_type}"

    # Add feature details text
    description = description.replace('\r', ' ')
    details_text = f"Theme: {theme}\nTable LDS: {table_lds}\nTable Type: {table_type}\n\nDescription:\n{description}"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.text = details_text
    # Force text to fit in the textbox
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)
    tf.word_wrap = True
    tf.auto_size = None

    add_image(slide, feature_picture_path, left=Inches(9), top=Inches(0), width=Inches(3.75), height=Inches(3.75))
    add_image(slide, aerial_picture_path, left=Inches(5), top=Inches(3.7), width=Inches(5), height=Inches(3.75))

if __name__ == "__main__":
    metadata_folder = r"C:\Data\model\metadata"
    metadata_file = os.path.join(metadata_folder, "metadata_info.json")
    # Read the metadata file
    with open(metadata_file, 'r') as f:
        metadata = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Create title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Topographic Data Model"

    summary_data = metadata['summary']
    metadata.pop('summary')

    for category_name, summary_info in summary_data.items():
        add_summary_slide(prs, summary_info)


    for table_key, table_val in metadata.items():
        layer_name = table_key
        feature_names = table_val['feature_types'] if 'feature_types' in table_val else []
        model_picture_path = table_val.get('model_picture_path', None)
        layer_description = table_val.get('layer_description', '')
        layer_category = table_val.get('layer_category', '')
        add_layer_slide(prs, layer_name, feature_names, layer_description, model_picture_path, layer_category=layer_category)

        # Add slides for each feature in the table
        if 'features' in table_val:
            for feature_name, feature_info in table_val['features'].items():
                feature_type = feature_name
                theme = feature_info.get('theme', '')
                table_lds = feature_info.get('table_lds', '')
                table_type = feature_info.get('table_type', '')
                description = feature_info.get('description', '')
                feature_picture_path = feature_info.get('feature_picture_path', None)
                aerial_picture_path = feature_info.get('aerial_picture_path', None)
                
                add_feature_slide(prs, feature_type, theme, table_lds, table_type, description, feature_picture_path, aerial_picture_path=aerial_picture_path)

    output_path = os.path.join(metadata_folder, "metadata_summary.pptx")
    prs.save(output_path)
    print(f"PowerPoint saved to {output_path}")
